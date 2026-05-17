# -*- coding: utf-8 -*-
"""
PPT Pipeline — Unified operations for beautification and creation.
DesignHub · 投影友好 · 一键美化 · 从零创建
"""
from __future__ import annotations

import argparse
import copy
import os
import sys
from dataclasses import dataclass, field
from typing import Optional

from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from .design import DesignSystem, load_theme, list_themes, FONT_SIZES


# ============================================================
# Utilities
# ============================================================

def _safe(fn, *args, **kwargs):
    """Call fn, return None on any exception."""
    try:
        return fn(*args, **kwargs)
    except Exception:
        return None


def _shape_has_text(shape) -> bool:
    """Check if shape contains actionable text."""
    return _safe(lambda: bool(shape.text_frame and shape.text.strip())) or False


def _is_placeholder(shape) -> bool:
    """Check if shape is a layout placeholder."""
    return _safe(lambda: shape.is_placeholder) or False


def _iter_text_runs(shape):
    """Yield all text runs across all paragraphs in a shape."""
    if not hasattr(shape, 'text_frame'):
        return
    tf = shape.text_frame
    if tf is None:
        return
    for para in tf.paragraphs:
        for run in para.runs:
            yield run


# ============================================================
# Slide type detection
# ============================================================

def detect_slide_type(slide) -> str:
    """Heuristic slide type detection.

    Returns one of:
        'cover'       — title slide (large text, few shapes)
        'section'     — section divider
        'content'     — standard content (text + optional images)
        'image'       — image-heavy (pictures > text)
        'comparison'  — two columns / side-by-side
        'ending'      — thank-you / ending slide
        'toc'         — table of contents / agenda
        'blank'       — mostly empty
        'unknown'     — fallback
    """
    shapes = list(slide.shapes)
    text_shapes = [s for s in shapes if _shape_has_text(s)]
    picture_shapes = [s for s in shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    text_count = len(text_shapes)
    pic_count = len(picture_shapes)

    # Slide-level text to check
    all_text = ' '.join(
        _safe(lambda: s.text_frame.text.strip()) or '' for s in text_shapes
    ).lower()

    # Ending slide
    if any(kw in all_text for kw in ['谢谢', 'thank', '感谢', 'q&a', '提问']):
        return 'ending'

    # TOC
    if any(kw in all_text for kw in ['目录', 'contents', 'agenda', 'outline', '大纲']):
        return 'toc'

    # Blank
    if len(shapes) == 0:
        return 'blank'

    # Cover: very few shapes, large text on title
    if text_count <= 3 and len(shapes) <= 5:
        font_sizes = []
        for s in text_shapes:
            for run in _iter_text_runs(s):
                if run.font.size:
                    font_sizes.append(run.font.size)
        if font_sizes and max(font_sizes) >= Pt(44):
            return 'cover'

    # Image-heavy: more pictures than text elements
    if pic_count >= 2 and pic_count >= text_count:
        return 'image'

    # Comparison: two distinct text columns with clear horizontal separation
    if text_count >= 4 and pic_count <= 2:
        positions = [
            _safe(lambda: s.left) for s in text_shapes
            if _shape_has_text(s)
        ]
        positions = sorted(p for p in positions if p is not None)
        if len(positions) >= 4:
            # Check for clear left-right clustering
            mid = (positions[0] + positions[-1]) / 2
            left_cluster = [p for p in positions if p < mid]
            right_cluster = [p for p in positions if p >= mid]
            if len(left_cluster) >= 2 and len(right_cluster) >= 2:
                gap = right_cluster[0] - left_cluster[-1]
                if gap > Inches(1.0):
                    return 'comparison'

    # Section divider
    if text_count <= 2 and pic_count <= 1 and len(shapes) <= 4:
        font_sizes = []
        for s in text_shapes:
            for run in _iter_text_runs(s):
                if run.font.size:
                    font_sizes.append(run.font.size)
        if font_sizes and max(font_sizes) >= Pt(38):
            return 'section'

    return 'content'


# ============================================================
# Beautification — Edit existing PPT
# ============================================================

def _apply_background(slide, color: RGBColor):
    """Set slide background to solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_left_accent_bar(slide, ds: DesignSystem):
    """Add vertical accent bar on left side."""
    color = ds.accent_bar_color or ds.primary
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        ds.left_bar_width, ds.slide_height,
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    # Move to bottom of z-order
    spTree = slide.shapes._spTree
    sp = bar._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def _add_top_accent_line(slide, ds: DesignSystem):
    """Add thin accent line at top."""
    color = ds.accent_bar_color or ds.primary
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        ds.slide_width, ds.top_line_height,
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def _add_page_number(slide, ds: DesignSystem, page_num: int):
    """Add page number at bottom-right."""
    left = ds.slide_width - Inches(1.0)
    top = ds.slide_height - Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, Inches(0.8), Inches(0.35))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(page_num)
    run.font.size = Pt(14)
    run.font.color.rgb = ds.text_dim
    run.font.name = ds.font_body


def _restyle_shape(shape, ds: DesignSystem, role: str = 'body'):
    """Apply theme styling to a text shape."""
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return  # never touch pictures

    for run in _iter_text_runs(shape):
        if role == 'title':
            run.font.name = ds.font_title
            run.font.color.rgb = ds.text_title
            if run.font.size is None or run.font.size < Pt(36):
                run.font.size = FONT_SIZES['page_title']
        elif role == 'subtitle':
            run.font.name = ds.font_body
            run.font.color.rgb = ds.text_body
            if run.font.size is None:
                run.font.size = FONT_SIZES['subtitle']
        elif role == 'section':
            run.font.name = ds.font_title
            run.font.color.rgb = ds.primary
            if run.font.size is None or run.font.size < Pt(36):
                run.font.size = FONT_SIZES['section_title']
        else:  # body
            run.font.name = ds.font_body
            run.font.color.rgb = ds.text_body
            if run.font.size is None:
                run.font.size = FONT_SIZES['body']

    # Style text frame
    if shape.has_text_frame:
        tf = shape.text_frame
        tf.word_wrap = True
        for para in tf.paragraphs:
            if para.alignment is None:
                para.alignment = PP_ALIGN.LEFT


def _style_slide_by_type(slide, slide_type: str, ds: DesignSystem, page_num: int):
    """Apply per-slide-type styling rules."""
    shapes = list(slide.shapes)
    text_shapes = [s for s in shapes if _shape_has_text(s)]

    if not text_shapes:
        return

    if slide_type in ('cover', 'ending'):
        # Center all text on cover/ending slides
        for s in text_shapes:
            if s.has_text_frame:
                for para in s.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.CENTER
            for run in _iter_text_runs(s):
                run.font.name = ds.font_title
                run.font.color.rgb = ds.text_title

    elif slide_type == 'section':
        # Section: big number + title, both centered
        for s in text_shapes:
            if s.has_text_frame:
                for para in s.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.CENTER
            _restyle_shape(s, ds, 'section')

    elif slide_type == 'content':
        # Content: first large shape = title, rest = body
        if len(text_shapes) >= 1:
            largest = max(
                text_shapes,
                key=lambda s: _safe(lambda: max(
                    r.font.size or Pt(0) for r in _iter_text_runs(s)
                )) or Pt(0)
            )
            for s in text_shapes:
                role = 'title' if s is largest else 'body'
                _restyle_shape(s, ds, role)

    else:
        # Default: style all as body
        for s in text_shapes:
            _restyle_shape(s, ds, 'body')


def beautify(
    input_path: str,
    theme: str = 'designhub_warm',
    output_path: Optional[str] = None,
    *,
    add_decor: bool = True,
    add_page_numbers: bool = True,
    preserve_images: bool = True,
    verbose: bool = True,
) -> str:
    """Beautify an existing PowerPoint presentation.

    Applies theme (colors, fonts, spacing), adds decorative elements,
    page numbers, and restyles all text while preserving images.

    Args:
        input_path: Path to input .pptx file.
        theme: Theme name ('designhub_warm', 'designhub_pro', 'projection').
        output_path: Output path (auto-generated if None).
        add_decor: Add left accent bar + top line + page numbers.
        add_page_numbers: Add page numbers to each slide.
        preserve_images: Keep all images untouched.
        verbose: Print progress.

    Returns:
        Path to the beautified .pptx file.
    """
    ds = load_theme(theme)
    prs = Presentation(input_path)

    if verbose:
        print(f'Beautifying: {input_path}')
        print(f'Theme: {ds.name}')
        print(f'Slides: {len(prs.slides)}')

    for i, slide in enumerate(prs.slides):
        page_num = i + 1

        # 1. Background
        _apply_background(slide, ds.bg)

        # 2. Decor
        if add_decor:
            _add_left_accent_bar(slide, ds)
            _add_top_accent_line(slide, ds)

        # 3. Page number
        if add_page_numbers:
            _add_page_number(slide, ds, page_num)

        # 4. Detect type and style
        slide_type = detect_slide_type(slide)
        _style_slide_by_type(slide, slide_type, ds, page_num)

        if verbose:
            print(f'  Slide {page_num:2d} → {slide_type}')

    # Output path
    if output_path is None:
        base, ext = os.path.splitext(input_path)
        output_path = f'{base}_{ds.name.replace(" · ", "_").replace(" ", "_")}{ext}'

    prs.save(output_path)
    if verbose:
        print(f'\nSaved: {output_path}')
    return output_path


# ============================================================
# Creation — Build new PPT from scratch
# ============================================================

@dataclass
class SlideSpec:
    """Specification for a single slide."""
    type: str = 'content'       # cover, section, content, image, ending, comparison, toc
    title: str = ''             # Slide title
    subtitle: str = ''          # Subtitle (cover/ending)
    body: list[str] = field(default_factory=list)  # Bullet points or body text
    images: list[str] = field(default_factory=list)  # Image file paths
    notes: str = ''             # Speaker notes
    layout: str = 'default'     # Layout template name


@dataclass
class OutlineSpec:
    """Complete presentation outline."""
    title: str = ''
    subtitle: str = ''
    author: str = ''
    date: str = ''
    slides: list[SlideSpec] = field(default_factory=list)


def _create_cover_slide(prs, spec: SlideSpec, ds: DesignSystem):
    """Create a cover/title slide."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    _apply_background(slide, ds.primary)

    # Title
    if spec.title:
        left = ds.margin_left
        top = Inches(2.0)
        width = ds.slide_width - ds.margin_left - ds.margin_right
        height = Inches(2.0)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = spec.title
        run.font.size = FONT_SIZES['hero']
        run.font.color.rgb = ds.text_light
        run.font.name = ds.font_title
        run.font.bold = True

    # Subtitle
    if spec.subtitle:
        top = Inches(4.2)
        txBox = slide.shapes.add_textbox(left, top, width, Inches(1.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = spec.subtitle
        run.font.size = FONT_SIZES['subtitle']
        run.font.color.rgb = ds.text_light
        run.font.name = ds.font_body

    # Bottom accent line
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        ds.margin_left, Inches(5.8),
        Inches(3.0), Inches(0.06),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ds.accent
    bar.line.fill.background()

    return slide


def _create_section_slide(prs, spec: SlideSpec, ds: DesignSystem):
    """Create a section divider slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _apply_background(slide, ds.bg)

    # Left accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(0.15), ds.slide_height,
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ds.primary
    bar.line.fill.background()

    # Section number (if available)
    title = spec.title
    number_str = ''
    parts = title.split('.', 1) if title else ['', '']
    if len(parts) == 2 and parts[0].isdigit():
        number_str, title = parts

    if number_str:
        txBox = slide.shapes.add_textbox(
            ds.margin_left + Inches(0.5),
            Inches(1.5),
            Inches(2.0), Inches(2.0),
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = number_str
        run.font.size = Pt(120)
        run.font.color.rgb = ds.primary
        run.font.name = ds.font_title
        run.font.bold = True
        run.font.color.rgb = RGBColor(
            ds.primary[0], ds.primary[1], ds.primary[2]
        )
        # 50% opacity via color (simplified: just use it as-is for now)

    # Section title
    txBox = slide.shapes.add_textbox(
        ds.margin_left + Inches(0.5),
        Inches(3.2),
        ds.slide_width - ds.margin_left - ds.margin_right,
        Inches(2.0),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title.strip()
    run.font.size = FONT_SIZES['section_title']
    run.font.color.rgb = ds.text_title
    run.font.name = ds.font_title

    return slide


def _create_content_slide(prs, spec: SlideSpec, ds: DesignSystem):
    """Create a standard content slide with title + bullets."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _apply_background(slide, ds.bg)

    # Left accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        ds.left_bar_width, ds.slide_height,
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ds.primary
    bar.line.fill.background()

    # Title
    if spec.title:
        title_left = ds.margin_left + Inches(0.2)
        txBox = slide.shapes.add_textbox(
            title_left, ds.margin_top,
            ds.slide_width - title_left - ds.margin_right,
            Inches(0.9),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = spec.title
        run.font.size = FONT_SIZES['page_title']
        run.font.color.rgb = ds.text_title
        run.font.name = ds.font_title
        run.font.bold = True

    # Title underline
    if spec.title:
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            ds.margin_left + Inches(0.2),
            ds.margin_top + Inches(0.9),
            Inches(2.5), Inches(0.05),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = ds.accent_bar_color or ds.primary
        line.line.fill.background()

    # Body bullets
    body_top = ds.margin_top + Inches(1.2)
    body_left = ds.margin_left + Inches(0.4)
    body_width = ds.slide_width - body_left - ds.margin_right

    txBox = slide.shapes.add_textbox(
        body_left, body_top,
        body_width, ds.slide_height - body_top - ds.margin_bottom,
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(spec.body):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(12)
        p.level = 0
        run = p.add_run()
        run.text = bullet
        run.font.size = FONT_SIZES['body']
        run.font.color.rgb = ds.text_body
        run.font.name = ds.font_body

    # Images (right side, if provided)
    if spec.images:
        img_count = len(spec.images)
        img_top = body_top
        img_width = Inches(4.5)
        img_left = ds.slide_width - img_width - ds.margin_right

        for i, img_path in enumerate(spec.images):
            if os.path.exists(img_path):
                img_height = Inches(2.5)
                slide.shapes.add_picture(
                    img_path,
                    img_left,
                    img_top + Inches(i * 2.8),
                    img_width,
                    img_height,
                )

    return slide


def _create_ending_slide(prs, spec: SlideSpec, ds: DesignSystem):
    """Create a thank-you ending slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _apply_background(slide, ds.primary)

    # Thank you text
    txBox = slide.shapes.add_textbox(
        Inches(1), Inches(2.5),
        ds.slide_width - Inches(2), Inches(2.0),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = spec.title or '谢谢！'
    run.font.size = FONT_SIZES['hero']
    run.font.color.rgb = ds.text_light
    run.font.name = ds.font_title
    run.font.bold = True

    if spec.subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(24)
        run2 = p2.add_run()
        run2.text = spec.subtitle
        run2.font.size = FONT_SIZES['body']
        run2.font.color.rgb = RGBColor(
            ds.text_light[0], ds.text_light[1], ds.text_light[2]
        )
        run2.font.name = ds.font_body

    return slide


_SLIDE_FACTORIES = {
    'cover': _create_cover_slide,
    'section': _create_section_slide,
    'content': _create_content_slide,
    'ending': _create_ending_slide,
}


def create_from_outline(
    outline: OutlineSpec,
    theme: str = 'designhub_warm',
    output_path: Optional[str] = None,
    *,
    verbose: bool = True,
) -> str:
    """Create a new presentation from a structured outline.

    Args:
        outline: OutlineSpec with title, subtitle, and list of SlideSpec.
        theme: Theme name.
        output_path: Output .pptx path (auto-generated if None).
        verbose: Print progress.

    Returns:
        Path to the created .pptx file.
    """
    ds = load_theme(theme)
    prs = Presentation()
    prs.slide_width = ds.slide_width
    prs.slide_height = ds.slide_height

    if verbose:
        print(f'Creating: {outline.title or "Untitled"}')
        print(f'Theme: {ds.name}')
        print(f'Slides: {len(outline.slides)}')

    for i, spec in enumerate(outline.slides):
        factory = _SLIDE_FACTORIES.get(spec.type, _create_content_slide)
        factory(prs, spec, ds)
        if verbose:
            print(f'  Slide {i + 1:2d} → {spec.type}: {spec.title or "(untitled)"}')

    if output_path is None:
        safe_title = (outline.title or 'presentation').replace(' ', '_').replace('/', '_')
        output_path = f'{safe_title}_{ds.name.replace(" · ", "_").replace(" ", "_")}.pptx'

    prs.save(output_path)
    if verbose:
        print(f'\nSaved: {output_path}')
    return output_path


# ============================================================
# Evaluation — PPTEval-inspired quality checklist
# ============================================================

@dataclass
class EvalResult:
    """Evaluation result for a single dimension."""
    score: int                # 1-5
    details: list[str] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)


@dataclass
class EvalReport:
    """Complete evaluation report for a presentation."""
    content: EvalResult
    design: EvalResult
    coherence: EvalResult
    file_path: str = ''

    @property
    def overall(self) -> float:
        return (self.content.score + self.design.score + self.coherence.score) / 3.0

    def summary(self) -> str:
        lines = [
            f'=== PPT Eval Report ===',
            f'File: {self.file_path}',
            f'',
            f'Content:   {self.content.score}/5',
            f'Design:    {self.design.score}/5',
            f'Coherence: {self.coherence.score}/5',
            f'Overall:   {self.overall:.1f}/5',
        ]
        for dim_name, result in [
            ('Content', self.content),
            ('Design', self.design),
            ('Coherence', self.coherence),
        ]:
            if result.warnings:
                lines.append(f'\n{dim_name} warnings:')
                for w in result.warnings:
                    lines.append(f'  [!] {w}')
        return '\n'.join(lines)


def evaluate(prs_path: str, *, verbose: bool = True) -> EvalReport:
    """Run PPTEval-inspired quality evaluation on a presentation.

    Evaluates three dimensions (1-5 scale):
      Content:   text clarity, completeness, readability
      Design:    color harmony, visual elements, consistency
      Coherence: logical flow, transitions, navigation

    Args:
        prs_path: Path to .pptx file.
        verbose: Print report.

    Returns:
        EvalReport with scores and warnings.
    """
    prs = Presentation(prs_path)

    # ===== Content evaluation =====
    content_score = 3  # default: average
    content_warnings = []

    total_slides = len(prs.slides)
    slides_with_text = 0
    empty_slides = 0
    text_lengths = []

    for slide in prs.slides:
        slide_text = ' '.join(
            _safe(lambda: s.text) or '' for s in slide.shapes if hasattr(s, 'text')
        )
        text_lengths.append(len(slide_text))
        if slide_text.strip():
            slides_with_text += 1
        else:
            empty_slides += 1

    text_ratio = slides_with_text / max(total_slides, 1)

    if text_ratio >= 0.9 and all(t > 0 for t in text_lengths):
        content_score = 4
    elif text_ratio >= 0.7:
        content_score = 3
    elif text_ratio >= 0.5:
        content_score = 2
        content_warnings.append(f'{empty_slides}/{total_slides} slides have no text')
    else:
        content_score = 1
        content_warnings.append(f'Most slides ({empty_slides}/{total_slides}) are empty')

    # ===== Design evaluation =====
    design_score = 3
    design_warnings = []

    bg_colors = set()
    for slide in prs.slides:
        bg = _safe(lambda: slide.background.fill.fore_color.rgb)
        if bg:
            bg_colors.add(str(bg))

    if len(bg_colors) == 0:
        design_score = 2
        design_warnings.append('No background colors set')
    elif len(bg_colors) == 1:
        design_score = 4  # consistent background
    elif len(bg_colors) <= 3:
        design_score = 3
    else:
        design_score = 2
        design_warnings.append(f'{len(bg_colors)} different background colors')

    # Count pictures per slide
    total_pics = sum(
        1 for slide in prs.slides
        for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE
    )
    if total_pics == 0 and total_slides > 3:
        design_warnings.append('No images or visual elements found')

    # ===== Coherence evaluation =====
    coherence_score = 3
    coherence_warnings = []

    has_title_slide = False
    has_ending_slide = False
    first_slide_text = ''
    last_slide_text = ''

    if prs.slides:
        first_slide_text = ' '.join(
            _safe(lambda: s.text) or '' for s in prs.slides[0].shapes if hasattr(s, 'text')
        ).lower()
        last_slide_text = ' '.join(
            _safe(lambda: s.text) or '' for s in prs.slides[-1].shapes if hasattr(s, 'text')
        ).lower()

    if any(kw in first_slide_text for kw in ['标题', 'title', '主题']):
        has_title_slide = True
    if any(kw in last_slide_text for kw in ['谢谢', 'thank', '感谢']):
        has_ending_slide = True

    if has_title_slide and has_ending_slide:
        coherence_score = 4
    elif has_title_slide or has_ending_slide:
        coherence_score = 3
        if not has_title_slide:
            coherence_warnings.append('No title/cover slide detected')
        if not has_ending_slide:
            coherence_warnings.append('No ending/thank-you slide detected')
    else:
        coherence_score = 2
        coherence_warnings.append('Missing both title and ending slides')

    report = EvalReport(
        content=EvalResult(score=content_score, warnings=content_warnings),
        design=EvalResult(score=design_score, warnings=design_warnings),
        coherence=EvalResult(score=coherence_score, warnings=coherence_warnings),
        file_path=prs_path,
    )

    if verbose:
        print(report.summary())
    return report


# ============================================================
# CLI — Command-line interface
# ============================================================

def cli_main(argv: Optional[list[str]] = None):
    """Command-line entry point."""
    parser = argparse.ArgumentParser(
        prog='ppt-pipeline',
        description='Unified PPT beautification and creation pipeline.',
    )
    sub = parser.add_subparsers(dest='command', required=True)

    # beautify
    p_beautify = sub.add_parser('beautify', help='Beautify an existing PPTX')
    p_beautify.add_argument('input', help='Input .pptx path')
    p_beautify.add_argument('-t', '--theme', default='designhub_warm',
                             help='Theme name (designhub_warm, designhub_pro, projection)')
    p_beautify.add_argument('-o', '--output', help='Output path')
    p_beautify.add_argument('--no-decor', action='store_true',
                             help='Skip decorative elements')
    p_beautify.add_argument('--no-numbers', action='store_true',
                             help='Skip page numbers')

    # create
    p_create = sub.add_parser('create', help='Create new PPTX from JSON outline')
    p_create.add_argument('outline', help='Path to outline JSON file')
    p_create.add_argument('-t', '--theme', default='designhub_warm',
                           help='Theme name')
    p_create.add_argument('-o', '--output', help='Output path')

    # evaluate
    p_eval = sub.add_parser('evaluate', help='Evaluate PPTX quality')
    p_eval.add_argument('input', help='Input .pptx path')

    # themes
    p_themes = sub.add_parser('themes', help='List available themes')

    args = parser.parse_args(argv)

    if args.command == 'themes':
        print('Available themes:')
        for t in list_themes():
            print(f'  - {t}')

    elif args.command == 'beautify':
        beautify(
            args.input,
            theme=args.theme,
            output_path=args.output,
            add_decor=not args.no_decor,
            add_page_numbers=not args.no_numbers,
        )

    elif args.command == 'create':
        import json
        with open(args.outline, 'r', encoding='utf-8') as f:
            data = json.load(f)
        outline = OutlineSpec(
            title=data.get('title', ''),
            subtitle=data.get('subtitle', ''),
            author=data.get('author', ''),
            date=data.get('date', ''),
            slides=[SlideSpec(**s) for s in data.get('slides', [])],
        )
        create_from_outline(outline, theme=args.theme, output_path=args.output)

    elif args.command == 'evaluate':
        evaluate(args.input)


if __name__ == '__main__':
    cli_main()
